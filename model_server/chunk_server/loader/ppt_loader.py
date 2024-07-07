from langchain_core.documents import Document
from pptx import Presentation
from tqdm import tqdm


class PPTLoader():
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        docs = []
        prs = Presentation(self.file_path)
        for slide_number, slide in tqdm(enumerate(prs.slides, start=1), desc=f"解析[{self.file_path}]的幻灯片"):
            full_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        full_text += paragraph.text.strip() + "\n"
                if shape.has_table:
                    table_content = ""
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                table_content += paragraph.text.strip() + "\n"
                    docs.append(Document(table_content, metadata={"format": "table"}))
            docs.append(Document(full_text.strip()))
        return docs
